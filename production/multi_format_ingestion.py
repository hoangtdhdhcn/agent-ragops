"""
Multi-format document ingestion pipeline for production RAG system.

Supports: PDF, HTML, DOCX, CSV, Excel, PowerPoint, Images (with OCR), and web scraping.
"""

import os
import re
import mimetypes
from pathlib import Path
from typing import List, Dict, Optional, Union, Tuple
from urllib.parse import urljoin, urlparse
import logging

import requests
from bs4 import BeautifulSoup
from PIL import Image
import pytesseract
import pandas as pd
from docx import Document as DocxDocument
from pptx import Presentation
from pdfminer.high_level import extract_text as extract_pdf_text
import pymupdf
import pymupdf4llm

from config import MARKDOWN_DIR

# Configure logging
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

class DocumentTypeDetector:
    """Detects document types and routes to appropriate processors."""
    
    SUPPORTED_FORMATS = {
        'pdf': ['application/pdf', '.pdf'],
        'html': ['text/html', '.html', '.htm'],
        'docx': ['application/vnd.openxmlformats-officedocument.wordprocessingml.document', '.docx'],
        'csv': ['text/csv', '.csv'],
        'excel': ['application/vnd.openxmlformats-officedocument.spreadsheetml.sheet', 
                 'application/vnd.ms-excel', '.xlsx', '.xls'],
        'pptx': ['application/vnd.openxmlformats-officedocument.presentationml.presentation', '.pptx'],
        'image': ['image/jpeg', 'image/png', 'image/tiff', 'image/bmp', 'image/gif', 
                 '.jpg', '.jpeg', '.png', '.tiff', '.tif', '.bmp', '.gif']
    }
    
    @classmethod
    def detect_format(cls, file_path: Union[str, Path], content_type: Optional[str] = None) -> Optional[str]:
        """Detect document format from file path and/or content type."""
        file_path = Path(file_path)
        file_extension = file_path.suffix.lower()
        
        # Check by file extension first
        for format_name, extensions in cls.SUPPORTED_FORMATS.items():
            if file_extension in extensions:
                return format_name
        
        # Check by content type if provided
        if content_type:
            for format_name, mime_types in cls.SUPPORTED_FORMATS.items():
                if content_type in mime_types:
                    return format_name
        
        return None
    
    @classmethod
    def is_supported(cls, file_path: Union[str, Path]) -> bool:
        """Check if file format is supported."""
        return cls.detect_format(file_path) is not None


class DocumentProcessor:
    """Base class for document processors."""
    
    def __init__(self, output_dir: Union[str, Path]):
        self.output_dir = Path(output_dir)
        self.output_dir.mkdir(parents=True, exist_ok=True)
    
    def process(self, file_path: Union[str, Path], metadata: Optional[Dict] = None) -> str:
        """Process document and return markdown content."""
        raise NotImplementedError


class PDFProcessor(DocumentProcessor):
    """Process PDF documents."""
    
    def process(self, file_path: Union[str, Path], metadata: Optional[Dict] = None) -> str:
        try:
            doc = pymupdf.open(file_path)
            md = pymupdf4llm.to_markdown(doc, header=False, footer=False, 
                                       page_separators=True, ignore_images=True, 
                                       write_images=False, image_path=None)
            md_cleaned = md.encode('utf-8', errors='surrogatepass').decode('utf-8', errors='ignore')
            return md_cleaned
        except Exception as e:
            logger.error(f"Error processing PDF {file_path}: {e}")
            return ""


class HTMLProcessor(DocumentProcessor):
    """Process HTML documents and web pages."""
    
    def process(self, file_path: Union[str, Path], metadata: Optional[Dict] = None) -> str:
        try:
            if self._is_url(file_path):
                return self._process_web_page(str(file_path), metadata)
            else:
                return self._process_local_html(file_path)
        except Exception as e:
            logger.error(f"Error processing HTML {file_path}: {e}")
            return ""
    
    def _is_url(self, path: str) -> bool:
        """Check if path is a URL."""
        try:
            result = urlparse(path)
            return all([result.scheme, result.netloc])
        except:
            return False
    
    def _process_web_page(self, url: str, metadata: Optional[Dict] = None) -> str:
        """Process web page content."""
        try:
            headers = {
                'User-Agent': 'Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/91.0.4472.124 Safari/537.36'
            }
            response = requests.get(url, headers=headers, timeout=30)
            response.raise_for_status()
            
            soup = BeautifulSoup(response.content, 'html.parser')
            
            # Remove unwanted elements
            for element in soup(['script', 'style', 'nav', 'header', 'footer', 'aside']):
                element.decompose()
            
            # Extract main content
            main_content = soup.find('main') or soup.find('article') or soup.find('div', class_=re.compile(r'main|content'))
            if main_content:
                text = main_content.get_text(separator='\n', strip=True)
            else:
                text = soup.get_text(separator='\n', strip=True)
            
            # Clean up text
            lines = [line.strip() for line in text.split('\n') if line.strip()]
            cleaned_text = '\n'.join(lines)
            
            # Add metadata
            md_content = f"# {metadata.get('title', urlparse(url).netloc)}\n\n"
            md_content += f"**Source:** {url}\n"
            md_content += f"**Retrieved:** {metadata.get('retrieved_date', '')}\n\n"
            md_content += cleaned_text
            
            return md_content
            
        except Exception as e:
            logger.error(f"Error processing web page {url}: {e}")
            return ""
    
    def _process_local_html(self, file_path: Union[str, Path]) -> str:
        """Process local HTML file."""
        try:
            with open(file_path, 'r', encoding='utf-8') as f:
                soup = BeautifulSoup(f, 'html.parser')
            
            # Remove unwanted elements
            for element in soup(['script', 'style']):
                element.decompose()
            
            text = soup.get_text(separator='\n', strip=True)
            lines = [line.strip() for line in text.split('\n') if line.strip()]
            return '\n'.join(lines)
            
        except Exception as e:
            logger.error(f"Error processing local HTML {file_path}: {e}")
            return ""


class DOCXProcessor(DocumentProcessor):
    """Process DOCX documents."""
    
    def process(self, file_path: Union[str, Path], metadata: Optional[Dict] = None) -> str:
        try:
            doc = DocxDocument(file_path)
            paragraphs = []
            
            # Extract text from paragraphs
            for para in doc.paragraphs:
                if para.text.strip():
                    paragraphs.append(para.text)
            
            # Extract text from tables
            for table in doc.tables:
                for row in table.rows:
                    row_text = []
                    for cell in row.cells:
                        cell_text = cell.text.strip()
                        if cell_text:
                            row_text.append(cell_text)
                    if row_text:
                        paragraphs.append(' | '.join(row_text))
            
            return '\n\n'.join(paragraphs)
            
        except Exception as e:
            logger.error(f"Error processing DOCX {file_path}: {e}")
            return ""


class CSVProcessor(DocumentProcessor):
    """Process CSV files."""
    
    def process(self, file_path: Union[str, Path], metadata: Optional[Dict] = None) -> str:
        try:
            df = pd.read_csv(file_path)
            return self._dataframe_to_markdown(df)
        except Exception as e:
            logger.error(f"Error processing CSV {file_path}: {e}")
            return ""
    
    def _dataframe_to_markdown(self, df: pd.DataFrame) -> str:
        """Convert DataFrame to markdown format."""
        if df.empty:
            return ""
        
        # Create markdown table
        md_lines = []
        
        # Header
        header = '| ' + ' | '.join(df.columns.astype(str)) + ' |'
        md_lines.append(header)
        
        # Separator
        separator = '| ' + ' | '.join(['---'] * len(df.columns)) + ' |'
        md_lines.append(separator)
        
        # Data rows
        for _, row in df.iterrows():
            row_str = '| ' + ' | '.join(row.astype(str)) + ' |'
            md_lines.append(row_str)
        
        return '\n'.join(md_lines)


class ExcelProcessor(DocumentProcessor):
    """Process Excel files (XLSX, XLS)."""
    
    def process(self, file_path: Union[str, Path], metadata: Optional[Dict] = None) -> str:
        try:
            # Read all sheets
            excel_file = pd.ExcelFile(file_path)
            md_content = []
            
            for sheet_name in excel_file.sheet_names:
                df = pd.read_excel(file_path, sheet_name=sheet_name)
                sheet_md = f"## Sheet: {sheet_name}\n\n"
                sheet_md += self._dataframe_to_markdown(df)
                md_content.append(sheet_md)
            
            return '\n\n'.join(md_content)
            
        except Exception as e:
            logger.error(f"Error processing Excel {file_path}: {e}")
            return ""
    
    def _dataframe_to_markdown(self, df: pd.DataFrame) -> str:
        """Convert DataFrame to markdown format."""
        if df.empty:
            return ""
        
        md_lines = []
        
        # Header
        header = '| ' + ' | '.join(df.columns.astype(str)) + ' |'
        md_lines.append(header)
        
        # Separator
        separator = '| ' + ' | '.join(['---'] * len(df.columns)) + ' |'
        md_lines.append(separator)
        
        # Data rows
        for _, row in df.iterrows():
            row_str = '| ' + ' | '.join(row.astype(str)) + ' |'
            md_lines.append(row_str)
        
        return '\n'.join(md_lines)


class PPTXProcessor(DocumentProcessor):
    """Process PowerPoint presentations."""
    
    def process(self, file_path: Union[str, Path], metadata: Optional[Dict] = None) -> str:
        try:
            presentation = Presentation(file_path)
            slides_content = []
            
            for i, slide in enumerate(presentation.slides):
                slide_text = []
                
                # Extract text from shapes
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text = shape.text.strip()
                        if text:
                            slide_text.append(text)
                
                if slide_text:
                    slide_md = f"## Slide {i + 1}\n\n"
                    slide_md += '\n\n'.join(slide_text)
                    slides_content.append(slide_md)
            
            return '\n\n'.join(slides_content)
            
        except Exception as e:
            logger.error(f"Error processing PPTX {file_path}: {e}")
            return ""


class ImageProcessor(DocumentProcessor):
    """Process images with OCR."""
    
    def __init__(self, output_dir: Union[str, Path]):
        super().__init__(output_dir)
        # Check if Tesseract is available
        try:
            pytesseract.get_tesseract_version()
        except:
            logger.warning("Tesseract OCR not found. Image processing will be limited.")
    
    def process(self, file_path: Union[str, Path], metadata: Optional[Dict] = None) -> str:
        try:
            # Try OCR first
            try:
                image = Image.open(file_path)
                text = pytesseract.image_to_string(image, lang='eng')
                if text.strip():
                    return f"# OCR Extracted Text from {Path(file_path).name}\n\n{text}"
            except:
                pass
            
            # Fallback: extract metadata and basic info
            try:
                image = Image.open(file_path)
                info = f"# Image: {Path(file_path).name}\n\n"
                info += f"**Format:** {image.format}\n"
                info += f"**Size:** {image.size}\n"
                info += f"**Mode:** {image.mode}\n"
                return info
            except:
                return f"# Image: {Path(file_path).name}\n\nUnable to process image file."
                
        except Exception as e:
            logger.error(f"Error processing image {file_path}: {e}")
            return ""


class MultiFormatIngestionPipeline:
    """Main pipeline for multi-format document ingestion."""
    
    def __init__(self, output_dir: Union[str, Path] = MARKDOWN_DIR):
        self.output_dir = Path(output_dir)
        self.output_dir.mkdir(parents=True, exist_ok=True)
        
        # Initialize processors
        self.processors = {
            'pdf': PDFProcessor(self.output_dir),
            'html': HTMLProcessor(self.output_dir),
            'docx': DOCXProcessor(self.output_dir),
            'csv': CSVProcessor(self.output_dir),
            'excel': ExcelProcessor(self.output_dir),
            'pptx': PPTXProcessor(self.output_dir),
            'image': ImageProcessor(self.output_dir),
        }
    
    def ingest_file(self, file_path: Union[str, Path], 
                   metadata: Optional[Dict] = None) -> Tuple[bool, str]:
        """Ingest a single file."""
        try:
            file_path = Path(file_path)
            
            if not file_path.exists():
                logger.error(f"File not found: {file_path}")
                return False, ""
            
            # Detect format
            format_type = DocumentTypeDetector.detect_format(file_path)
            if not format_type:
                logger.warning(f"Unsupported format for: {file_path}")
                return False, ""
            
            # Process document
            processor = self.processors[format_type]
            markdown_content = processor.process(file_path, metadata)
            
            if not markdown_content.strip():
                logger.warning(f"No content extracted from: {file_path}")
                return False, ""
            
            # Save markdown file
            output_path = self.output_dir / f"{file_path.stem}.md"
            output_path.write_text(markdown_content, encoding='utf-8')
            
            logger.info(f"Successfully processed: {file_path} -> {output_path}")
            return True, str(output_path)
            
        except Exception as e:
            logger.error(f"Error ingesting file {file_path}: {e}")
            return False, ""
    
    def ingest_directory(self, directory: Union[str, Path], 
                        recursive: bool = True) -> Dict[str, int]:
        """Ingest all supported files in a directory."""
        directory = Path(directory)
        if not directory.exists():
            logger.error(f"Directory not found: {directory}")
            return {}
        
        stats = {'processed': 0, 'failed': 0, 'skipped': 0}
        
        if recursive:
            files = [f for f in directory.rglob('*') if f.is_file()]
        else:
            files = [f for f in directory.iterdir() if f.is_file()]
        
        for file_path in files:
            if DocumentTypeDetector.is_supported(file_path):
                success, _ = self.ingest_file(file_path)
                if success:
                    stats['processed'] += 1
                else:
                    stats['failed'] += 1
            else:
                stats['skipped'] += 1
        
        logger.info(f"Ingestion completed. Stats: {stats}")
        return stats
    
    def ingest_web_page(self, url: str, metadata: Optional[Dict] = None) -> Tuple[bool, str]:
        """Ingest a web page."""
        try:
            processor = self.processors['html']
            markdown_content = processor.process(url, metadata)
            
            if not markdown_content.strip():
                return False, ""
            
            # Generate filename from URL
            parsed_url = urlparse(url)
            filename = f"web_{parsed_url.netloc}_{parsed_url.path.replace('/', '_') or 'index'}.md"
            output_path = self.output_dir / filename
            output_path.write_text(markdown_content, encoding='utf-8')
            
            logger.info(f"Successfully processed web page: {url} -> {output_path}")
            return True, str(output_path)
            
        except Exception as e:
            logger.error(f"Error ingesting web page {url}: {e}")
            return False, ""


# Utility functions for easy use
def ingest_single_file(file_path: Union[str, Path], 
                      output_dir: Union[str, Path] = MARKDOWN_DIR) -> Tuple[bool, str]:
    """Convenience function to ingest a single file."""
    pipeline = MultiFormatIngestionPipeline(output_dir)
    return pipeline.ingest_file(file_path)


def ingest_directory(directory: Union[str, Path], 
                    output_dir: Union[str, Path] = MARKDOWN_DIR,
                    recursive: bool = True) -> Dict[str, int]:
    """Convenience function to ingest a directory."""
    pipeline = MultiFormatIngestionPipeline(output_dir)
    return pipeline.ingest_directory(directory, recursive)


def ingest_web_page(url: str, 
                   output_dir: Union[str, Path] = MARKDOWN_DIR,
                   metadata: Optional[Dict] = None) -> Tuple[bool, str]:
    """Convenience function to ingest a web page."""
    pipeline = MultiFormatIngestionPipeline(output_dir)
    return pipeline.ingest_web_page(url, metadata)


if __name__ == "__main__":
    
    pipeline = MultiFormatIngestionPipeline()
    
    # Ingest a single file
    # success, path = pipeline.ingest_file("path/to/document.pdf")
    
    # Ingest a directory
    # stats = pipeline.ingest_directory("path/to/documents/")
    
    # Ingest a web page
    # success, path = pipeline.ingest_web_page("https://example.com")
    pass